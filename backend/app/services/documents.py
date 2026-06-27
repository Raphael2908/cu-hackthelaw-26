from __future__ import annotations

import io
from pathlib import Path

# Best-effort plain-text extraction from an uploaded file. The demo attaches text to a case
# (architecture.md §9) — no object store. PDF/DOCX/PPTX use lightweight pure-Python parsers;
# everything else is decoded as UTF-8. There is no OCR, so scanned/image PDFs and image-only
# slide decks yield no text and are rejected.

SUPPORTED = {".pdf", ".docx", ".pptx", ".txt", ".md", ".markdown", ".text"}


def _extract_pdf(raw: bytes) -> str:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(raw))
    return "\n\n".join((page.extract_text() or "") for page in reader.pages).strip()


def _extract_docx(raw: bytes) -> str:
    import docx  # python-docx

    document = docx.Document(io.BytesIO(raw))
    return "\n".join(p.text for p in document.paragraphs).strip()


def _extract_pptx(raw: bytes) -> str:
    from pptx import Presentation  # python-pptx

    presentation = Presentation(io.BytesIO(raw))
    slides = []
    for slide in presentation.slides:
        texts = [
            shape.text for shape in slide.shapes if shape.has_text_frame and shape.text.strip()
        ]
        if texts:
            slides.append("\n".join(texts))
    return "\n\n".join(slides).strip()


def extract_text(filename: str, raw: bytes) -> str:
    """Return the plain text of an uploaded document, or raise ValueError if the format is
    unsupported or no text could be extracted (e.g. an empty or scanned-image PDF)."""
    suffix = Path(filename or "").suffix.lower()
    if suffix not in SUPPORTED and suffix != "":
        raise ValueError(f"Unsupported file type '{suffix or filename}'.")

    try:
        if suffix == ".pdf":
            text = _extract_pdf(raw)
        elif suffix == ".docx":
            text = _extract_docx(raw)
        elif suffix == ".pptx":
            text = _extract_pptx(raw)
        else:
            text = raw.decode("utf-8", errors="ignore").strip()
    except ValueError:
        raise
    except Exception as e:  # corrupt / encrypted / unparseable file → clean error, never a 500
        raise ValueError(f"Could not parse '{filename}': {e}") from e

    if not text:
        raise ValueError(f"No text could be extracted from '{filename}'.")
    return text
