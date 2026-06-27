from __future__ import annotations

import io

import pytest

from app.services import documents


def _build_pptx(*texts: str) -> bytes:
    """Build a minimal in-memory .pptx with one title-only slide per text. Offline, no fixtures."""
    from pptx import Presentation

    presentation = Presentation()
    blank = presentation.slide_layouts[6]  # fully blank layout
    title = presentation.slide_layouts[5]  # "Title Only"
    for text in texts:
        slide = presentation.slides.add_slide(title if text else blank)
        if text:
            slide.shapes.title.text = text
    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def test_pptx_in_supported():
    assert ".pptx" in documents.SUPPORTED


def test_extract_pptx_returns_slide_text():
    raw = _build_pptx("Indemnity overview", "Governing law: England")
    text = documents.extract_text("deck.pptx", raw)
    assert "Indemnity overview" in text
    assert "Governing law: England" in text


def test_extract_pptx_empty_deck_raises():
    raw = _build_pptx("")  # one blank slide, no text → 415 path
    with pytest.raises(ValueError):
        documents.extract_text("blank.pptx", raw)
